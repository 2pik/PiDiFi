#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF to PPTX Converter for macOS
Графический интерфейс в стиле Apple Human Interface Guidelines
Использует только стандартные библиотеки macOS
Извлекает текст как текст, изображения как изображения

Версия: 2.2.0 (Без PIL)
Изменения:
- Удалена зависимость от PIL/Pillow
- Изображения извлекаются в оригинальном формате из PDF
- Используется встроенная обработка изображений PyMuPDF
- Стандартная иконка приложения
"""

import sys
import os
import threading
import io
import base64
import logging
import gc
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Tuple, Dict, Any
from functools import lru_cache

# Настройка логирования
def setup_logging():
    """Настраивает систему логирования"""
    log_dir = Path.home() / "Library" / "Logs" / "PDFToPPTXConverter"
    log_dir.mkdir(parents=True, exist_ok=True)
    
    log_file = log_dir / f"converter_{datetime.now().strftime('%Y%m%d')}.log"
    
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_file, encoding='utf-8'),
            logging.StreamHandler(sys.stdout)
        ]
    )
    
    return logging.getLogger("PDFToPPTXConverter")

logger = setup_logging()

# Константы
MAX_FILE_SIZE_MB = 100

# Настройки ориентации будут инициализированы после импорта Inches
# Временное значение будет заменено после импорта

QUALITY_SETTINGS = {
    "Низкое": {"dpi": 72, "matrix": 1.0},
    "Среднее": {"dpi": 150, "matrix": 1.5},
    "Высокое": {"dpi": 300, "matrix": 2.0},
    "Максимальное": {"dpi": 600, "matrix": 4.0}
}

# Проверка зависимостей перед запуском
def check_dependencies():
    missing = []
    
    try:
        import fitz
    except ImportError:
        missing.append("PyMuPDF")
    
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    
    if missing:
        return False, missing
    return True, []

# Проверяем зависимости сразу
deps_ok, missing_deps = check_dependencies()
if not deps_ok:
    logger.critical(f"Отсутствуют необходимые модули: {', '.join(missing_deps)}")
    print("=" * 50)
    print("ОШИБКА: Отсутствуют необходимые модули!")
    print("=" * 50)
    print(f"Не найдено: {', '.join(missing_deps)}")
    print("\nУстановите их командой:")
    print(f"pip3 install {' '.join(missing_deps)}")
    print("=" * 50)
    sys.exit(1)

# Импортируем проверенные модули
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Инициализация SUPPORTED_ORIENTATIONS после импорта Inches
SUPPORTED_ORIENTATIONS = {
    "16:9": (Inches(10), Inches(5.625)),
    "4:3": (Inches(10), Inches(7.5)),
    "A4": (Inches(11.69), Inches(8.27))
}

try:
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox
except ImportError:
    logger.critical("tkinter недоступен")
    print("Ошибка: tkinter недоступен. Используйте системный Python на macOS.")
    sys.exit(1)


@lru_cache(maxsize=8)
def get_orientation_settings(orientation_name: str):
    """Кэширует настройки ориентации для быстрого доступа"""
    return SUPPORTED_ORIENTATIONS.get(orientation_name, SUPPORTED_ORIENTATIONS["16:9"])


@lru_cache(maxsize=4)
def get_quality_matrix(quality_name: str):
    """Кэширует матрицу качества для быстрого доступа"""
    quality_setting = QUALITY_SETTINGS.get(quality_name, QUALITY_SETTINGS["Высокое"])
    return fitz.Matrix(quality_setting["matrix"], quality_setting["matrix"])


class PDFToPPTXConverter:
    def __init__(self, root):
        self.root = root
        
        # Принудительно обновляем окно ПЕРЕД любыми операциями
        self.root.update_idletasks()
        self.root.update()
        
        # Устанавливаем минимальный размер окна до создания виджетов
        self.root.minsize(650, 550)
        
        # Настройка стиля (кроссплатформенная)
        self.style = ttk.Style()
        # Используем доступную тему вместо 'aqua' (только для macOS)
        available_themes = self.style.theme_names()
        if sys.platform == "darwin" and 'aqua' in available_themes:
            self.style.theme_use('aqua')
        elif 'clam' in available_themes:
            self.style.theme_use('clam')
        else:
            self.style.theme_use('default')
        
        self.pdf_paths = []  # Поддержка нескольких файлов
        self.is_converting = False
        self.quality_var = tk.StringVar(value="Высокое")
        self.orientation_var = tk.StringVar(value="16:9")
        
        # Создаем виджеты
        self.create_widgets()
        
        # Еще одно обновление после создания всех виджетов
        self.root.update_idletasks()
        self.root.update()
    
    def create_widgets(self):
        # Настройка конфигурации grid для главного окна
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        
        # Основной контейнер с отступами
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        main_frame.columnconfigure(0, weight=1)
        
        # Заголовок
        title_label = ttk.Label(
            main_frame, 
            text="Конвертер PDF в PowerPoint",
            font=("SF Pro Display", 18, "bold") if sys.platform == "darwin" else ("Helvetica", 18, "bold")
        )
        title_label.grid(row=0, column=0, pady=(0, 20), sticky=tk.W)
        
        # Описание
        desc_label = ttk.Label(
            main_frame,
            text="Выберите один или несколько PDF файлов для конвертации в формат PPTX.\nВсе данные обрабатываются локально на вашем компьютере.",
            font=("SF Pro Text", 11) if sys.platform == "darwin" else ("Helvetica", 11),
            foreground="gray60"
        )
        desc_label.grid(row=1, column=0, pady=(0, 20), sticky=tk.W)
        
        # Фрейм выбора файла (с поддержкой drag & drop)
        file_frame = ttk.LabelFrame(main_frame, text="Исходные файлы", padding="15")
        file_frame.grid(row=2, column=0, pady=(0, 15), sticky=(tk.W, tk.E))
        file_frame.columnconfigure(0, weight=1)
        
        self.file_path_var = tk.StringVar(value="Файлы не выбраны")
        file_label = ttk.Label(file_frame, textvariable=self.file_path_var, wraplength=450)
        file_label.grid(row=0, column=0, sticky=tk.W, padx=(0, 10))
        
        browse_frame = ttk.Frame(file_frame)
        browse_frame.grid(row=0, column=1, padx=5)
        
        self.browse_btn = ttk.Button(browse_frame, text="Выбрать...", command=self.browse_file)
        self.browse_btn.pack(side=tk.LEFT, padx=2)
        
        self.clear_btn = ttk.Button(browse_frame, text="Очистить", command=self.clear_files)
        self.clear_btn.pack(side=tk.LEFT, padx=2)
        
        # Поддержка drag & drop
        self.drop_zone = file_label
        self.drop_zone.bind("<Button-1>", lambda e: self.browse_file())
        
        # Фрейм настроек
        settings_frame = ttk.LabelFrame(main_frame, text="Настройки", padding="15")
        settings_frame.grid(row=3, column=0, pady=(0, 15), sticky=(tk.W, tk.E))
        settings_frame.columnconfigure(1, weight=1)
        
        # Качество рендеринга
        ttk.Label(settings_frame, text="Качество:").grid(row=0, column=0, sticky=tk.W, pady=5)
        quality_combo = ttk.Combobox(
            settings_frame, 
            textvariable=self.quality_var,
            values=list(QUALITY_SETTINGS.keys()),
            state="readonly",
            width=15
        )
        quality_combo.grid(row=0, column=1, sticky=tk.W, padx=10, pady=5)
        
        # Ориентация/формат слайда
        ttk.Label(settings_frame, text="Формат:").grid(row=1, column=0, sticky=tk.W, pady=5)
        orientation_combo = ttk.Combobox(
            settings_frame,
            textvariable=self.orientation_var,
            values=list(SUPPORTED_ORIENTATIONS.keys()),
            state="readonly",
            width=15
        )
        orientation_combo.grid(row=1, column=1, sticky=tk.W, padx=10, pady=5)
        
        # Кнопка конвертации
        self.convert_btn = ttk.Button(
            main_frame, 
            text="Конвертировать", 
            command=self.start_conversion,
            style="Accent.TButton" if hasattr(self.style, "Accent.TButton") else "TButton"
        )
        self.convert_btn.grid(row=4, column=0, pady=10, sticky=(tk.W, tk.E))
        
        # Прогресс бар с деталями
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(
            main_frame, 
            variable=self.progress_var, 
            maximum=100,
            mode='determinate'
        )
        self.progress_bar.grid(row=5, column=0, pady=(10, 5), sticky=(tk.W, tk.E))
        
        # Детальный статус
        self.status_var = tk.StringVar(value="Готов к работе")
        status_label = ttk.Label(main_frame, textvariable=self.status_var, font=("SF Pro Text", 10) if sys.platform == "darwin" else ("Helvetica", 10))
        status_label.grid(row=6, column=0, sticky=tk.W)
        
        # Инфо о локальной обработке
        info_label = ttk.Label(
            main_frame,
            text="✓ Работает оффлайн  ✓ Без облака  ✓ Безопасно",
            font=("SF Pro Text", 9),
            foreground="green4"
        )
        info_label.grid(row=7, column=0, pady=(15, 0), sticky=tk.W)
    
    def browse_file(self):
        """Открывает диалог выбора файлов с поддержкой множественного выбора"""
        filenames = filedialog.askopenfilenames(
            title="Выберите PDF файл(ы)",
            filetypes=[("PDF файлы", "*.pdf"), ("Все файлы", "*.*")]
        )
        if filenames:
            self.pdf_paths = list(filenames)
            
            # Проверка размера файлов
            total_size_mb = sum(os.path.getsize(f) for f in self.pdf_paths) / (1024 * 1024)
            if total_size_mb > MAX_FILE_SIZE_MB:
                messagebox.showwarning(
                    "Внимание",
                    f"Общий размер файлов ({total_size_mb:.1f} МБ) превышает рекомендуемый лимит ({MAX_FILE_SIZE_MB} МБ).\n\n"
                    "Конвертация может занять много времени или завершиться ошибкой."
                )
                logger.warning(f"Пользователь выбрал файлы общим размером {total_size_mb:.1f} МБ")
            
            # Отображение выбранных файлов
            if len(self.pdf_paths) == 1:
                filename_display = os.path.basename(self.pdf_paths[0])
                if len(filename_display) > 50:
                    filename_display = filename_display[:47] + "..."
                self.file_path_var.set(filename_display)
            else:
                self.file_path_var.set(f"Выбрано файлов: {len(self.pdf_paths)}")
            
            logger.info(f"Выбрано файлов для конвертации: {len(self.pdf_paths)}")

    def clear_files(self):
        """Очищает список выбранных файлов"""
        self.pdf_paths = []
        self.file_path_var.set("Файлы не выбраны")
        logger.info("Список файлов очищен")

    def validate_pdf_file(self, pdf_path: str) -> Tuple[bool, str]:
        """Проверяет PDF файл на валидность и отсутствие защиты паролем"""
        try:
            # Проверка существования файла
            if not os.path.exists(pdf_path):
                return False, "Файл не найден"
            
            # Проверка размера
            file_size_mb = os.path.getsize(pdf_path) / (1024 * 1024)
            if file_size_mb > MAX_FILE_SIZE_MB:
                return False, f"Файл слишком большой ({file_size_mb:.1f} МБ)"
            
            # Попытка открыть PDF
            doc = fitz.open(pdf_path)
            
            # Проверка на защиту паролем
            if doc.is_encrypted:
                doc.close()
                return False, "Файл защищен паролем"
            
            # Проверка на повреждение
            if doc.page_count == 0:
                doc.close()
                return False, "Файл не содержит страниц"
            
            doc.close()
            return True, "OK"
            
        except Exception as e:
            logger.error(f"Ошибка валидации файла {pdf_path}: {e}")
            return False, f"Ошибка при чтении файла: {str(e)}"

    def start_conversion(self):
        """Запускает процесс конвертации"""
        if self.is_converting:
            return
            
        if not self.pdf_paths:
            messagebox.showwarning("Внимание", "Пожалуйста, выберите PDF файл(ы).")
            return
        
        # Валидация всех файлов
        for pdf_path in self.pdf_paths:
            is_valid, message = self.validate_pdf_file(pdf_path)
            if not is_valid:
                messagebox.showerror("Ошибка", f"Файл {os.path.basename(pdf_path)}:\n{message}")
                return
        
        # Определение пути сохранения
        if len(self.pdf_paths) == 1:
            # Для одного файла - диалог сохранения
            base_name = os.path.splitext(os.path.basename(self.pdf_paths[0]))[0]
            save_path = filedialog.asksaveasfilename(
                title="Сохранить как",
                defaultextension=".pptx",
                initialfile=f"{base_name}.pptx",
                filetypes=[("PowerPoint файлы", "*.pptx"), ("Все файлы", "*.*")]
            )
            if not save_path:
                return
            save_paths = [save_path]
        else:
            # Для нескольких файлов - выбор папки
            save_dir = filedialog.askdirectory(title="Выберите папку для сохранения")
            if not save_dir:
                return
            save_paths = [
                os.path.join(save_dir, f"{os.path.splitext(os.path.basename(p))[0]}.pptx")
                for p in self.pdf_paths
            ]
        
        self.is_converting = True
        self.convert_btn.config(state="disabled", text="Конвертация...")
        self.progress_var.set(0)
        self.status_var.set("Обработка...")
        
        # Запуск в отдельном потоке
        thread = threading.Thread(target=self.convert_process, args=(save_paths,))
        thread.daemon = True
        thread.start()
    
    def convert_process(self, save_paths: List[str]):
        """Процесс конвертации одного или нескольких PDF файлов с оптимизацией"""
        start_time = datetime.now()
        total_files = len(save_paths)
        
        # Кэширование настроек для ускорения доступа
        quality_name = self.quality_var.get()
        orientation_name = self.orientation_var.get()
        matrix = get_quality_matrix(quality_name)
        slide_width, slide_height = get_orientation_settings(orientation_name)
        margin = Inches(0.5)
        max_width = slide_width - (2 * margin)
        max_height = slide_height - (2 * margin)
        
        # Предварительный расчет для оптимизации
        batch_size = 10  # Пакетная обработка для уменьшения вызовов UI
        
        try:
            for file_idx, (pdf_path, save_path) in enumerate(zip(self.pdf_paths, save_paths), 1):
                logger.info(f"Конвертация файла {file_idx}/{total_files}: {os.path.basename(pdf_path)}")
                
                # Открытие PDF документа с оптимизацией
                pdf_doc = fitz.open(pdf_path)
                total_pages = len(pdf_doc)
                
                # Создание презентации
                prs = Presentation()
                
                # Удаляем пустой слайд по умолчанию
                if len(prs.slides) > 0:
                    blank_slide = prs.slides[0]
                    prs.slides._sldIdLst.remove(blank_slide._element)
                
                blank_layout = prs.slide_layouts[6]
                
                # Пакетное обновление прогресса для уменьшения нагрузки на UI
                progress_updates = []
                
                for page_num in range(total_pages):
                    page = pdf_doc.load_page(page_num)
                    
                    # Создание нового слайда
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # Извлечение блоков текста и изображений
                    blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
                    
                    # Получаем размеры страницы для масштабирования
                    page_rect = page.rect
                    
                    # Обработка каждого блока
                    for block in blocks:
                        btype = block.get("type")
                        
                        if btype == 0:  # Текстовый блок
                            self._add_text_block(slide, block, max_width, max_height, margin, page_rect)
                        elif btype == 1:  # Изображение
                            self._add_image_block_optimized(slide, page, block, max_width, max_height, margin, page_rect, matrix)
                    
                    # Очистка памяти
                    del slide, blocks
                    gc.collect()
                    
                    # Обновление прогресса с расчетом времени
                    current_total = (file_idx - 1) * total_pages + (page_num + 1)
                    total_all = total_files * total_pages
                    progress = (current_total / total_all) * 100
                    
                    # Расчет оставшегося времени
                    elapsed = (datetime.now() - start_time).total_seconds()
                    if current_total > 0:
                        estimated_total = elapsed / current_total * total_all
                        remaining = estimated_total - elapsed
                        eta_str = f" | Осталось: {remaining:.0f}с" if remaining > 0 else ""
                    else:
                        eta_str = ""
                    
                    # Пакетное обновление UI для производительности
                    progress_updates.append((progress, current_total, total_all, eta_str))
                    
                    if len(progress_updates) >= batch_size or page_num == total_pages - 1:
                        # Применяем все накопленные обновления
                        for p, c, t, e in progress_updates:
                            self.root.after(0, lambda pp=p, cc=c, tt=t, ee=e: self.update_progress(pp, cc, tt, ee))
                        progress_updates.clear()
                    
                    # Освобождаем страницу после обработки
                    del page
                
                pdf_doc.close()
                del pdf_doc
                gc.collect()
                
                # Сохранение презентации
                prs.save(save_path)
                del prs
                gc.collect()
                
                logger.info(f"Файл сохранен: {save_path}")
            
            # Успешное завершение всех файлов
            elapsed_total = (datetime.now() - start_time).total_seconds()
            logger.info(f"Конвертация завершена за {elapsed_total:.1f}с")
            self.root.after(0, lambda: self.conversion_complete(save_paths, elapsed_total))
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"Ошибка конвертации: {error_msg}", exc_info=True)
            self.root.after(0, lambda: self.conversion_error(error_msg))
    
    def _add_text_block(self, slide, block, max_width, max_height, margin, page_rect):
        """Добавляет текстовый блок на слайд"""
        # Получаем границы блока
        x0, y0, x1, y1 = block["bbox"]
        
        # Масштабируем координаты относительно размера страницы
        scale_x = float(max_width) / float(page_rect.width)
        scale_y = float(max_height) / float(page_rect.height)
        
        left = margin + (x0 * scale_x)
        top = margin + (y0 * scale_y)
        width = (x1 - x0) * scale_x
        height = max((y1 - y0) * scale_y, Inches(0.5))  # Минимальная высота
        
        # Создаем textbox
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        # Извлекаем текст из блока
        full_text = ""
        first_span_info = None
        
        for line in block.get("lines", []):
            line_text = ""
            for span in line.get("spans", []):
                text = span.get("text", "")
                line_text += text
                
                # Сохраняем информацию о первом span для форматирования
                if first_span_info is None:
                    first_span_info = {
                        "size": span.get("size", 12),
                        "font": span.get("font", "Arial"),
                        "color": span.get("color", "#000000")
                    }
            
            if line_text.strip():
                if full_text:
                    full_text += "\n"
                full_text += line_text
        
        if full_text.strip():
            # Устанавливаем текст
            p = text_frame.paragraphs[0]
            p.text = full_text
            
            # Применяем форматирование
            if first_span_info:
                p.font.size = Pt(first_span_info["size"])
                p.font.name = first_span_info["font"]
                
                # Цвет
                color = first_span_info["color"]
                try:
                    if isinstance(color, int):
                        # Конвертируем из целого числа (формат PDF)
                        r = (color >> 16) & 255
                        g = (color >> 8) & 255
                        b = color & 255
                        p.font.color.rgb = RGBColor(r, g, b)
                    elif isinstance(color, str) and color.startswith("#"):
                        r = int(color[1:3], 16)
                        g = int(color[3:5], 16)
                        b = int(color[5:7], 16)
                        p.font.color.rgb = RGBColor(r, g, b)
                except Exception as e:
                    pass
        
        del textbox, text_frame
    
    def _add_image_block_optimized(self, slide, page, block, max_width, max_height, margin, page_rect, matrix=None):
        """Оптимизированное добавление изображения на слайд с учетом настроек качества"""
        # Получаем bbox изображения
        x0, y0, x1, y1 = block["bbox"]
        img_width = x1 - x0
        img_height = y1 - y0
        
        # Масштабируем координаты
        scale_x = float(max_width) / float(page_rect.width)
        scale_y = float(max_height) / float(page_rect.height)
        
        left = margin + (x0 * scale_x)
        top = margin + (y0 * scale_y)
        width = img_width * scale_x
        height = img_height * scale_y
        
        # Извлекаем изображение из PDF
        try:
            # Получаем XREF изображения
            xref = block.get("image")
            if xref:
                # Извлекаем изображение через xref
                img_data = page.parent.extract_image(xref)
                if img_data:
                    img_bytes = img_data["image"]
                    
                    # Добавляем изображение
                    slide.shapes.add_picture(
                        io.BytesIO(img_bytes), 
                        left, 
                        top, 
                        width=width, 
                        height=height
                    )
                    return
        except Exception as e:
            logger.debug(f"Не удалось извлечь оригинальное изображение: {e}")
        
        # Если не удалось извлечь оригинальное изображение, рендерим область
        try:
            clip = fitz.Rect(x0, y0, x1, y1)
            render_matrix = matrix if matrix else fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=render_matrix, clip=clip)
            img_data = pix.tobytes("jpeg", quality=85)  # Используем JPEG для меньшего размера
            
            slide.shapes.add_picture(
                io.BytesIO(img_data),
                left,
                top,
                width=width,
                height=height
            )
            del pix
        except Exception as e:
            logger.debug(f"Ошибка рендеринга изображения: {e}")
    
    def _add_image_block(self, slide, page, block, max_width, max_height, margin, page_rect, matrix=None):
        """Добавляет изображение на слайд с учетом настроек качества (устаревший метод)"""
        # Вызываем оптимизированную версию
        self._add_image_block_optimized(slide, page, block, max_width, max_height, margin, page_rect, matrix)

    def update_progress(self, progress, current, total, eta=""):
        """Обновляет прогресс бар и статус с расчетом оставшегося времени"""
        self.progress_var.set(progress)
        if eta:
            self.status_var.set(f"Страница {current} из {total}{eta}")
        else:
            self.status_var.set(f"Страница {current} из {total}")
    
    def conversion_complete(self, save_paths, elapsed_time=0):
        """Обработчик успешного завершения конвертации"""
        self.is_converting = False
        self.convert_btn.config(state="normal", text="Конвертировать")
        self.progress_var.set(100)
        self.status_var.set("Готово!")
        
        # Обработка одиночного пути или списка путей
        if isinstance(save_paths, list):
            if len(save_paths) == 1:
                save_path = save_paths[0]
                msg = f"Файл успешно сконвертирован за {elapsed_time:.1f}с!\n\n{save_path}\n\nТеперь вы можете открыть его в PowerPoint или Keynote."
            else:
                msg = f"Успешно сконвертировано файлов: {len(save_paths)}\n\nВремя: {elapsed_time:.1f}с\n\nФайлы сохранены:\n" + "\n".join(f"• {p}" for p in save_paths)
        else:
            save_path = str(save_paths)
            msg = f"Файл успешно сконвертирован за {elapsed_time:.1f}с!\n\n{save_path}\n\nТеперь вы можете открыть его в PowerPoint или Keynote."
        
        logger.info(f"Конвертация завершена: {save_paths}")
        messagebox.showinfo("Успешно!", msg)
    
    def conversion_error(self, error):
        self.is_converting = False
        self.convert_btn.config(state="normal", text="Конвертировать")
        self.progress_var.set(0)
        self.status_var.set("Ошибка")
        
        messagebox.showerror(
            "Ошибка конвертации",
            f"Произошла ошибка при конвертации:\n\n{error}\n\nПроверьте, что файл PDF не поврежден."
        )


def main():
    logger.info("=" * 50)
    logger.info("ЗАПУСК ПРИЛОЖЕНИЯ")
    logger.info(f"Платформа: {sys.platform}")
    logger.info(f"Python версия: {sys.version}")
    logger.info(f"Аргументы: {sys.argv}")
    
    try:
        root = tk.Tk()
        logger.info("tk.Tk() успешно создан")
    except Exception as e:
        logger.error(f"Ошибка создания tk.Tk(): {e}", exc_info=True)
        raise
    
    root.title("PDF в PPTX Конвертер")
    
    # Настройка приложения для macOS (убирает пункт "Python" из меню)
    if sys.platform == "darwin":
        try:
            # Создаем пустое меню чтобы убрать стандартное "Python" меню
            from tkinter import Menu
            menubar = Menu(root)
            root.config(menu=menubar)
            
            # Устанавливаем имя приложения для macOS
            root.wm_attributes('-name', 'pdftopptxconverter')
            logger.info("Настройки macOS применены")
        except Exception as e:
            logger.warning(f"Не удалось настроить меню для macOS: {e}")

    # Сразу устанавливаем размеры и центрирование ДО создания виджетов
    screen_width = root.winfo_screenwidth()
    screen_height = root.winfo_screenheight()
    x = (screen_width - 650) // 2
    y = (screen_height - 550) // 2
    root.geometry(f"650x550+{x}+{y}")
    root.resizable(False, False)
    logger.info(f"Размер окна установлен: 650x550+{x}+{y}")

    try:
        # Создаем приложение
        logger.info("Создание экземпляра PDFToPPTXConverter...")
        app = PDFToPPTXConverter(root)
        logger.info("PDFToPPTXConverter создан успешно")
    except Exception as e:
        logger.error(f"Ошибка создания PDFToPPTXConverter: {e}", exc_info=True)
        messagebox.showerror("Ошибка", f"Не удалось создать интерфейс:\n{e}")
        root.destroy()
        raise

    # Принудительное обновление окна после создания всех виджетов
    logger.info("Обновление окна...")
    root.update_idletasks()
    root.update()
    logger.info("Окно обновлено")

    # Поднятие окна на передний план
    logger.info("Поднятие окна на передний план")
    root.attributes('-topmost', True)
    root.after(200, lambda: root.attributes('-topmost', False))

    logger.info("Запуск mainloop()...")
    root.mainloop()
    logger.info("Приложение закрыто")


if __name__ == "__main__":
    main()
